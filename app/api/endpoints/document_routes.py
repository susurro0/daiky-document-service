import os
from datetime import datetime
from transformers import pipeline, AutoTokenizer
from fastapi import APIRouter, HTTPException, UploadFile, File
from starlette.responses import JSONResponse

from app.api.schemas.parsed_document_schema import ParsedDocument
from app.api.schemas.document_schemas import DocumentCreate, Document
from app.crud.document_crud import DocumentCRUD
from app.dependencies import Dependency

class DocumentRoutes:
    def __init__(self, dependency: Dependency, document_crud=DocumentCRUD):
        self.router = APIRouter()
        self.db = dependency.get_db()
        self.document_crud = DocumentCRUD(db=self.db)
        self.tokenizer = AutoTokenizer.from_pretrained("facebook/bart-large-cnn")
        self.summarizer = pipeline("summarization")

        @self.router.post("/api/upload/", response_model=Document)
        def upload_file(file: UploadFile = File(...)):
            """
            Upload a file and save its details in the database.
            """
            try:
                file_name = file.filename
                file_type = file.content_type
                upload_timestamp = datetime.now()
                file_location = f"uploads/{file_name}"

                if os.path.exists(file_location):
                    print(f"File {file_name} already exists. Overwriting...")

                with open(file_location, "wb") as f:
                    f.write(file.file.read())

                document_create = DocumentCreate(
                    file=file,
                    file_name=file_name,
                    file_type=file_type,
                    upload_timestamp=upload_timestamp,
                )

                saved_document = self.document_crud.create_document(document_create)

                return JSONResponse(
                    content={
                        "id": saved_document.id,
                        "file_name": saved_document.file_name,
                        "file_type": saved_document.file_type,
                        "upload_timestamp": saved_document.upload_timestamp.isoformat(),
                        "parsed_text": saved_document.parsed_text,
                    },
                    status_code=200,
                )

            except Exception as e:
                print(f"Failed to upload file: {e}")
                raise HTTPException(
                    status_code=500, detail="An error occurred while uploading the file."
                )

        def __split_with_overlap(text, tokenizer_name="bert-base-uncased", max_length=256, overlap=50):
            """
            Split text into overlapping chunks using the tokenizer.
            """
            tokens = self.tokenizer.encode(text, add_special_tokens=False)
            chunks = [
                self.tokenizer.decode(tokens[i:i + max_length], skip_special_tokens=True)
                for i in range(0, len(tokens), max_length - overlap)
            ]
            return chunks

        @self.router.get("/api/documents/{document_id}/parse", response_model=ParsedDocument)
        def parse_document(document_id: int):
            """
            Parse the content of a document and summarize it.
            """
            try:
                document = self.document_crud.get_document(document_id=document_id)
                if document is None:
                    raise HTTPException(status_code=404, detail="Document not found")

                file_location = f"uploads/{document.file_name}"
                text = ""

                if document.file_type == "PDF":
                    from pypdf import PdfReader
                    with open(file_location, "rb") as f:
                        pdf_reader = PdfReader(f)
                        text = "".join(page.extract_text() for page in pdf_reader.pages)

                elif document.file_type == "DOCX":
                    from docx import Document as DocxDocument
                    docx = DocxDocument(file_location)
                    text = "\n".join(para.text for para in docx.paragraphs)

                elif document.file_type == "PPTX":
                    from pptx import Presentation
                    pptx = Presentation(file_location)
                    text = "\n".join(
                        shape.text for slide in pptx.slides for shape in slide.shapes if hasattr(shape, "text")
                    )
                else:
                    raise HTTPException(status_code=415, detail="Unsupported file format")

                chunks = __split_with_overlap(text, max_length=128, overlap=25)
                summary = __summarize_text(text)

                return ParsedDocument(chunks=chunks, summary=summary)

            except Exception as e:
                if isinstance(e, HTTPException):
                    raise e
                print(f"Failed to parse document: {e}")
                raise HTTPException(
                    status_code=500,
                    detail="An error occurred while parsing the document."
                )

        def __summarize_text(text, max_length=512, summary_max_length=150, summary_min_length=25):
            """
            Summarize text using a pre-trained model.
            """
            tokens = self.tokenizer(text, return_tensors="pt", truncation=False)
            if len(tokens["input_ids"][0]) >= max_length:
                summary = self.summarizer(
                    text, max_length=summary_max_length, min_length=summary_min_length, do_sample=False
                )
                return summary[0]["summary_text"]
            return ""

